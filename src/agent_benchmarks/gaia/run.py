"""
GAIA runner for the Lightpanda agent (web-browsing subset).

Loads GAIA from HuggingFace (config `2023_level1`, split `validation`),
keeps only tasks with no attached file (the pure web-browsing subset),
invokes `lightpanda agent --task` one-shot per row, and grades answers
with the normalized exact-match rubric from the GAIA paper.

The GAIA dataset is gated on HuggingFace — accept the terms at
https://huggingface.co/datasets/gaia-benchmark/GAIA and set HF_TOKEN.

Per GAIA's license, predictions.jsonl contains the question and gold
answer verbatim, so results/ MUST remain gitignored (it is, at the
benchmarks/ root .gitignore).

Example:

    HF_TOKEN=... uv run gaia-run --limit 3
    HF_TOKEN=... uv run gaia-run --workers 8
"""

from __future__ import annotations

import argparse
import contextlib
import os
import sys
import tempfile
from pathlib import Path
from typing import Any

from datasets import load_dataset  # type: ignore[import-not-found]

from ..common import (
    add_common_runner_args,
    emit_scores,
    extract_answer_envelope,
    load_completed_ids,
    print_lightpanda_missing,
    resolve_lightpanda_binary,
    resolve_out_dir,
    run_benchmark_tasks,
    run_lightpanda_task,
    write_run_manifest,
)
from .grade import grade_predictions

OFFICE_EXTS = {".docx", ".xlsx", ".pptx"}


def _extract_office_text(path: Path) -> str:
    """Extract a plain-text representation of an office document. Returns a
    format that preserves as much structure as the model can use: tables
    emit one row per line with tab-separated cells; pptx emits per-slide
    headings. Visual-only information (cell colors, shapes) is lost —
    tasks depending on it will still score 0."""
    ext = path.suffix.lower()
    if ext == ".docx":
        from docx import Document  # type: ignore[import-not-found]

        doc = Document(str(path))
        parts: list[str] = [p.text for p in doc.paragraphs if p.text]
        for table in doc.tables:
            for row in table.rows:
                parts.append("\t".join(cell.text for cell in row.cells))
        return "\n".join(parts)
    if ext == ".xlsx":
        from openpyxl import load_workbook  # type: ignore[import-not-found]

        wb = load_workbook(str(path), data_only=True)
        out: list[str] = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            out.append(f"=== Sheet: {sheet_name} ===")
            for row in sheet.iter_rows(values_only=True):
                out.append("\t".join("" if c is None else str(c) for c in row))
        return "\n".join(out)
    if ext == ".pptx":
        from pptx import Presentation  # type: ignore[import-not-found]

        prs = Presentation(str(path))
        out: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            out.append(f"=== Slide {i} ===")
            for shape in slide.shapes:
                text = getattr(shape, "text", None)
                if text:
                    out.append(text)
        return "\n".join(out)
    raise ValueError(f"unknown office extension: {ext}")


def _preprocess_attachment(path: Path) -> Path:
    """Normalize an attachment for Lightpanda. Office docs are extracted to
    plain text written to a sibling .txt tempfile so Lightpanda's text-file
    handler picks them up. Other formats pass through untouched."""
    if path.suffix.lower() not in OFFICE_EXTS:
        return path
    text = _extract_office_text(path)
    fd, tmp = tempfile.mkstemp(suffix=f"__{path.name}.txt", text=True)
    try:
        with os.fdopen(fd, "w", encoding="utf-8") as f:
            f.write(text)
    except Exception:
        os.unlink(tmp)
        raise
    return Path(tmp)


# benchmarks/src/agent_benchmarks/gaia/run.py → benchmarks/
PROJECT_ROOT = Path(__file__).resolve().parents[3]

SYSTEM_PROMPT = """\
You are a research assistant driving the Lightpanda headless browser on the GAIA QA benchmark.

The Lightpanda browser tools are the ONLY way you can access the web. There is no WebSearch, no WebFetch, no shortcut — you must navigate real pages. Your tool surface includes `search`, `goto`, `tree`, `markdown`, `extract`, `structuredData`, `findElement`, `interactiveElements`, `links`, `click`, `fill`, `hover`, `selectOption`, `setChecked`, `press`, `scroll`, `waitForSelector`, `nodeDetails`, `getUrl`, `eval`, `consoleLogs`, `detectForms`.

BE PERSISTENT — this is the load-bearing instruction:
- GAIA tasks expect multi-step browsing. Most need 20-50 tool calls; some need 100+. Answers from prior knowledge without browsing score 0.
- If a search returns poor results, try DIFFERENT phrasings — synonyms, narrower queries, different angles.
- If a page is unreachable, find a DIFFERENT source. Wikipedia, official sites, archived pages, news outlets.
- If extraction fails, try a different tool (markdown → tree → extract → structuredData → findElement).
- Do NOT respond "unknown" or fall back to prior knowledge until you have made at least 20 substantive tool calls AND tried at least 3 different sources/angles.
- Small-candidate questions ("A, B, or C", yes/no): always pick one — never abstain.

Strategy:
1. Plan: prefer authoritative direct sources (Wikipedia, official sites) over search-engine landing pages when you know where to go.
2. Search: use the `search` tool — do NOT goto google.com directly. With `TAVILY_API_KEY` set, `search` queries Tavily and returns a clean numbered list of {title, url, snippet}; without the key, it falls back to scraping the DuckDuckGo HTML endpoint. Google scraping is blocked by Lightpanda's User-Agent and TLS fingerprint.
3. Navigate the available browser tools. Re-inspect after page-changing actions — DOM snapshots and node ids go stale.
4. Cross-check on a second source where the answer is non-obvious.

Final-answer envelope — STRICT
================================
Your entire response will be discarded except for the LAST text wrapped in `<ANSWER>...</ANSWER>` tags. Reasoning, tool-call narration, partial-credit notes — anything outside the envelope — is ignored. Only the envelope contents are graded.

GAIA grades by exact match after normalization (lowercase, strip articles/punct).

Format INSIDE the envelope:
- No preface, no explanation, no markdown, no source citations.
- Numbers: bare digits — no comma separators, no currency unless asked, no units beyond what's asked.
- Names/titles: complete and verbatim, no decoration, no surrounding punctuation.
- Lists: comma-separated on one line.

Examples (good):
  <ANSWER>45</ANSWER>
  <ANSWER>Oko, Thief of Crowns</ANSWER>
  <ANSWER>Paris, London, Tokyo</ANSWER>

Examples (bad — these score 0):
  <ANSWER>**45**</ANSWER>                              ← markdown inside
  <ANSWER>The answer is 45.</ANSWER>                   ← preface inside
  <ANSWER>$45</ANSWER>                                 ← currency
  <ANSWER>1,234</ANSWER>                               ← comma separator

Tool-use rules:
- Never use backendNodeId with click, fill, hover, selectOption, or setChecked. Always use a CSS selector.
- Use findElement to resolve a description into a selector when needed.
- Use distinguishing attributes (value, name, position) so selectors are unique.
- For credentials, pass $LP_USERNAME / $LP_PASSWORD directly as values.
"""

TASK_PROMPT_TEMPLATE = (
    "{task}\n\n"
    "Wrap your final answer in <ANSWER>...</ANSWER>. Inside the envelope: "
    "bare digits for numbers, comma-separated on one line for lists, no markdown."
)


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    add_common_runner_args(parser, suite_name="gaia")
    parser.add_argument(
        "--level",
        type=int,
        default=1,
        choices=[1, 2, 3],
        help="GAIA level (default: 1 — shortest tool-use chains, web-browsing focused)",
    )
    parser.add_argument(
        "--skip-attachments",
        action="store_true",
        help="Skip rows with attached files. Default is to include them, even though "
        "Lightpanda can't read PDFs/audio/images — skipped tasks score 0 and the full "
        "Level-N score is what the GAIA paper reports. Use this flag only when you "
        "want to isolate agent performance on text-only tasks.",
    )
    args = parser.parse_args(argv)

    lightpanda = resolve_lightpanda_binary(args.lightpanda, PROJECT_ROOT)
    if not lightpanda.exists():
        print_lightpanda_missing(lightpanda)
        return 2

    out_dir = resolve_out_dir(args.out_dir, PROJECT_ROOT, "gaia")
    predictions_path = out_dir / "predictions.jsonl"
    write_run_manifest(out_dir, agent_provider=args.provider, agent_model=args.model)

    completed = load_completed_ids(predictions_path) if args.resume else set()
    if completed:
        print(f"Resuming — skipping {len(completed)} already-completed task(s)", file=sys.stderr)

    config = f"2023_level{args.level}"
    print(f"Loading gaia-benchmark/GAIA {config}/{args.split} from HuggingFace...", file=sys.stderr)
    if not os.environ.get("HF_TOKEN") and not os.environ.get("HUGGING_FACE_HUB_TOKEN"):
        print(
            "warning: no HF_TOKEN in env — GAIA is a gated dataset and this load will fail.\n"
            "  1. Accept terms at https://huggingface.co/datasets/gaia-benchmark/GAIA\n"
            "  2. export HF_TOKEN=hf_...",
            file=sys.stderr,
        )
    ds = load_dataset("gaia-benchmark/GAIA", config, split=args.split)
    rows: list[dict[str, Any]] = list(ds)

    # Materialize the dataset snapshot locally so we can point `--attach`
    # at real file paths. Only needed when attachments aren't skipped.
    snapshot_dir: Path | None = None
    if not args.skip_attachments and any((r.get("file_name") or "").strip() for r in rows):
        from huggingface_hub import snapshot_download  # type: ignore[import-not-found]

        print("Downloading GAIA dataset snapshot for attachments...", file=sys.stderr)
        snapshot_dir = Path(snapshot_download(repo_id="gaia-benchmark/GAIA", repo_type="dataset"))

    if args.skip_attachments:
        before = len(rows)
        rows = [r for r in rows if not (r.get("file_name") or "").strip()]
        print(
            f"Skipped attachments: {before - len(rows)}/{before} rows dropped (kept {len(rows)})",
            file=sys.stderr,
        )

    if args.limit is not None:
        rows = rows[: args.limit]

    pending = [r for r in rows if r["task_id"] not in completed]

    def _work(row: dict[str, Any]) -> dict[str, Any]:
        attachment: Path | None = None
        tmp_to_delete: Path | None = None
        file_name = (row.get("file_name") or "").strip()
        if file_name and snapshot_dir is not None:
            raw = snapshot_dir / row["file_path"]
            if raw.exists():
                try:
                    attachment = _preprocess_attachment(raw)
                    if attachment != raw:
                        tmp_to_delete = attachment
                except Exception as e:
                    print(
                        f"warn: preprocess failed for {raw}: {e}; passing raw path",
                        file=sys.stderr,
                    )
                    attachment = raw
        try:
            pred, duration_s, timed_out, stderr_tail, rc, trace, usage = run_lightpanda_task(
                lightpanda=lightpanda,
                provider=args.provider,
                model=args.model,
                user_agent=args.user_agent,
                system_prompt=SYSTEM_PROMPT,
                task_prompt=TASK_PROMPT_TEMPLATE.format(task=row["Question"]),
                attachment=attachment,
                timeout_s=args.timeout,
            )
        finally:
            if tmp_to_delete is not None:
                with contextlib.suppress(OSError):
                    tmp_to_delete.unlink()
        pred, envelope_note = extract_answer_envelope(pred)
        if envelope_note:
            stderr_tail = (stderr_tail or "") + f"\n[envelope]\n{envelope_note}\n"
        return {
            "id": row["task_id"],
            "task": row["Question"],
            "gold": row["Final answer"],
            "prediction": pred,
            "duration_s": round(duration_s, 2),
            "timed_out": timed_out,
            "returncode": rc,
            "level": row.get("Level"),
            "file_name": file_name or None,
            "trace": trace,
            "usage": usage,
            "stderr_tail": stderr_tail,
        }

    run_benchmark_tasks(
        pending,
        _work,
        predictions_path=predictions_path,
        workers=args.workers,
        timeout_s=args.timeout,
        provider=args.provider,
        model=args.model,
        preview_fn=lambda row: row["Question"],
    )

    print("\nGrading...", file=sys.stderr)
    emit_scores(grade_predictions(predictions_path), out_dir)
    return 0


if __name__ == "__main__":
    sys.exit(main())
